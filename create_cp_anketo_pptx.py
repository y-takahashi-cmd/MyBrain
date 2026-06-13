import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

matplotlib.rcParams['font.family'] = 'Meiryo'

PRIMARY   = '#74B2E0'
SECONDARY = '#1E3A5F'
LIGHT     = '#B8D9F0'
WHITE     = '#FFFFFF'

def rgb(hex_color):
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if line:
        s.line.color.rgb = rgb(line)
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, l, t, w, h, size=12, bold=False, color=SECONDARY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return tb

def header(slide, title):
    add_rect(slide, 0, 0, 10, 0.9, SECONDARY)
    add_text(slide, title, 0.3, 0.1, 9.4, 0.7, size=16, bold=True, color=WHITE)
    add_rect(slide, 0, 0, 10, 7.5, WHITE)  # 白背景（後ろに）
    # 白背景を後ろに送る
    sp = slide.shapes[-1]._element
    slide.shapes[0]._element.addprevious(sp)

def header_slide(slide, title):
    add_rect(slide, 0, 0, 10, 7.5, WHITE)
    add_rect(slide, 0, 0, 10, 0.9, SECONDARY)
    add_text(slide, title, 0.3, 0.1, 9.4, 0.7, size=16, bold=True, color=WHITE)

def make_pie():
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_facecolor('white')
    labels = ['元気に働いている\n3件（60%）', '少し疲れている\n2件（40%）']
    ax.pie([60,40], labels=labels, colors=[PRIMARY, LIGHT], startangle=90,
           textprops={'fontsize':11,'color':SECONDARY})
    ax.axis('equal')
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    buf.seek(0)
    return buf

def make_bar(items, values):
    n = len(items)
    fig, ax = plt.subplots(figsize=(6.5, max(2.5, n*0.55)))
    fig.patch.set_facecolor('white')
    ir = items[::-1]
    vr = values[::-1]
    y  = np.arange(n)
    bars = ax.barh(y, vr, color=PRIMARY, height=0.5)
    ax.set_yticks(y)
    ax.set_yticklabels(ir, fontsize=10)
    ax.set_xlim(0, 110)
    ax.set_xlabel('%', fontsize=9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    for bar, val in zip(bars, vr):
        if val > 0:
            ax.text(val+1, bar.get_y()+bar.get_height()/2,
                    f'{val}%', va='center', fontsize=10, color=SECONDARY)
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    buf.seek(0)
    return buf

# ===== 表紙 =====
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, 10, 7.5, WHITE)
add_rect(s1, 0.5, 1.2, 9, 0.05, PRIMARY)
add_rect(s1, 0.5, 6.2, 9, 0.05, PRIMARY)
add_text(s1, 'クラウドパワー株式会社',           1, 1.5, 8, 0.8, size=22, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER)
add_text(s1, '企業のオンライン保健室',             1, 2.4, 8, 0.6, size=14, color=SECONDARY,  align=PP_ALIGN.CENTER)
add_text(s1, '社員アンケート集計結果',             1, 3.1, 8, 0.9, size=30, bold=True, color=PRIMARY,   align=PP_ALIGN.CENTER)
add_text(s1, '2026年6月',                         1, 4.1, 8, 0.6, size=14, color=SECONDARY,  align=PP_ALIGN.CENTER)
add_text(s1, '回答数：5件（全員回答）',            1, 4.7, 8, 0.6, size=13, color='#555555',  align=PP_ALIGN.CENTER)
add_text(s1, '企業のオンライン保健室　髙橋由香',   1, 5.4, 8, 0.6, size=12, color=SECONDARY,  align=PP_ALIGN.CENTER)

# ===== 全体所見 =====
s2 = prs.slides.add_slide(blank)
add_rect(s2, 0, 0, 10, 7.5, WHITE)
add_rect(s2, 0, 0, 10, 0.9, SECONDARY)
add_text(s2, '全体所見', 0.3, 0.1, 9.4, 0.7, size=20, bold=True, color=WHITE)
insights = [
    ('01', '表面上の状態は安定',   '元気に働いている60%。ただし40%が疲れを自覚している。'),
    ('02', '身体面の潜在的な疲れ', '睡眠・疲労感・ストレスが同率40%。自覚症状の薄い層も存在。'),
    ('03', '相談文化が薄い',       '自分で解決する80%・専門家0%。保健室がまだ使われていない。'),
    ('04', '保健室への期待はある', '健康相談・ストレス相談が各60%。ニーズ自体は存在する。'),
    ('05', '要確認事項',           'Q3「仕事の進め方」60%：背景（効率化・習得・不満）の確認が必要。'),
]
for i,(num,ttl,body) in enumerate(insights):
    top = 1.1 + i*1.1
    add_rect(s2, 0.4, top, 0.55, 0.65, PRIMARY)
    add_text(s2, num,  0.4, top+0.1, 0.55, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, ttl,  1.1, top,     4.5,  0.4, size=13, bold=True, color=SECONDARY)
    add_text(s2, body, 1.1, top+0.42, 8.5, 0.55, size=10, color='#444444')

# ===== Q1 =====
s3 = prs.slides.add_slide(blank)
add_rect(s3, 0, 0, 10, 7.5, WHITE)
add_rect(s3, 0, 0, 10, 0.9, SECONDARY)
add_text(s3, 'Q1　現在の状態に最も近いものを選んでください', 0.3, 0.1, 9.4, 0.7, size=15, bold=True, color=WHITE)
add_text(s3, '単一選択・回答 5件', 0.3, 1.0, 9, 0.4, size=11, color='#666666')
s3.shapes.add_picture(make_pie(), Inches(2.5), Inches(1.5), Inches(5), Inches(4.2))

# ===== Q2 =====
s4 = prs.slides.add_slide(blank)
add_rect(s4, 0, 0, 10, 7.5, WHITE)
add_rect(s4, 0, 0, 10, 0.9, SECONDARY)
add_text(s4, 'Q2　現在、気になっていることを教えてください', 0.3, 0.1, 9.4, 0.7, size=15, bold=True, color=WHITE)
add_text(s4, '複数回答・回答 5件', 0.3, 1.0, 9, 0.4, size=11, color='#666666')
s4.shapes.add_picture(
    make_bar(['睡眠','疲労感','ストレス','将来への不安','お金','特になし'],[40,40,40,20,20,40]),
    Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
add_text(s4, '睡眠・疲労・ストレスが同率40%。身体と精神の両面に潜在的な疲れ。', 0.5, 6.2, 9, 0.6, size=11, color=SECONDARY)

# ===== Q3 =====
s5 = prs.slides.add_slide(blank)
add_rect(s5, 0, 0, 10, 7.5, WHITE)
add_rect(s5, 0, 0, 10, 0.9, SECONDARY)
add_text(s5, 'Q3　今後、変えていきたいことは何ですか', 0.3, 0.1, 9.4, 0.7, size=15, bold=True, color=WHITE)
add_text(s5, '複数回答・回答 5件', 0.3, 1.0, 9, 0.4, size=11, color='#666666')
s5.shapes.add_picture(
    make_bar(['仕事の進め方','睡眠','ストレス対策','人間関係','コミュニケーション','時間管理'],[60,40,40,20,20,20]),
    Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
add_text(s5, '「仕事の進め方」60%トップ。背景（効率化・やり方習得・環境不満）は要確認。', 0.5, 6.2, 9, 0.6, size=11, color=SECONDARY)

# ===== Q4 =====
s6 = prs.slides.add_slide(blank)
add_rect(s6, 0, 0, 10, 7.5, WHITE)
add_rect(s6, 0, 0, 10, 0.9, SECONDARY)
add_text(s6, 'Q4　困った時に相談する相手を教えてください', 0.3, 0.1, 9.4, 0.7, size=15, bold=True, color=WHITE)
add_text(s6, '複数回答・回答 5件', 0.3, 1.0, 9, 0.4, size=11, color='#666666')
s6.shapes.add_picture(
    make_bar(['自分で解決する','友人','上司','家族','同僚','専門家'],[80,60,60,40,40,0]),
    Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
add_text(s6, '「自分で解決」80%・「専門家」0%。保健室がまだ相談の場として認識されていない。', 0.5, 6.2, 9, 0.6, size=11, color=SECONDARY)

# ===== Q5 =====
s7 = prs.slides.add_slide(blank)
add_rect(s7, 0, 0, 10, 7.5, WHITE)
add_rect(s7, 0, 0, 10, 0.9, SECONDARY)
add_text(s7, 'Q5　会社にあったら嬉しいサポートを教えてください', 0.3, 0.1, 9.4, 0.7, size=15, bold=True, color=WHITE)
add_text(s7, '複数回答・回答 5件', 0.3, 1.0, 9, 0.4, size=11, color='#666666')
s7.shapes.add_picture(
    make_bar(['仕事の整理や考え方','特に必要ない','メンタル相談','個別相談'],[40,40,20,20]),
    Inches(1.5), Inches(1.5), Inches(7), Inches(4))
add_text(s7, '「仕事の整理」と「特に必要ない」が同率40%。押しつけにならない関わり方が重要。', 0.5, 6.2, 9, 0.6, size=11, color=SECONDARY)

# ===== Q6 =====
s8 = prs.slides.add_slide(blank)
add_rect(s8, 0, 0, 10, 7.5, WHITE)
add_rect(s8, 0, 0, 10, 0.9, SECONDARY)
add_text(s8, 'Q6　オンライン保健室に今後期待することを教えてください', 0.3, 0.1, 9.4, 0.7, size=15, bold=True, color=WHITE)
add_text(s8, '複数回答・回答 5件', 0.3, 1.0, 9, 0.4, size=11, color='#666666')
s8.shapes.add_picture(
    make_bar(['健康相談','ストレス相談','よく分からない'],[60,60,20]),
    Inches(1.5), Inches(1.5), Inches(7), Inches(3.5))
add_text(s8, '健康・ストレス相談への期待が各60%。「よく分からない」20%→保健室の役割説明が途上。', 0.5, 6.2, 9, 0.6, size=11, color=SECONDARY)

# ===== Q7 ＋ 総括 =====
s9 = prs.slides.add_slide(blank)
add_rect(s9, 0, 0, 10, 7.5, WHITE)
add_rect(s9, 0, 0, 10, 0.9, SECONDARY)
add_text(s9, 'Q7　自由記述 ／ 総括', 0.3, 0.1, 9.4, 0.7, size=18, bold=True, color=WHITE)

add_rect(s9, 0.5, 1.1, 9, 1.7, LIGHT)
add_text(s9, 'Q7　自由記述（1名の声）', 0.7, 1.15, 8, 0.45, size=12, bold=True, color=SECONDARY)
add_text(s9, '「社長の意見が絶対になってしまうため、\n社長のいない会話の場を設ける必要があると考えている」',
         0.7, 1.6, 8.5, 1.0, size=12, color=SECONDARY)

add_text(s9, '総括', 0.5, 3.0, 9, 0.45, size=14, bold=True, color=SECONDARY)
for i, item in enumerate([
    '・保健室への期待はあるが、「自分で解決する」文化が根強く相談のハードルが高い状態',
    '・Q3「仕事の進め方」の背景確認が今後の関わり方の鍵',
    '・「特に必要ない」40%を踏まえ、押しつけにならない関わり方の設計が重要',
    '・Q7は1名の個人意見として参考情報に留める',
]):
    add_text(s9, item, 0.5, 3.55 + i*0.65, 9, 0.6, size=11, color='#333333')

out = r'C:\Users\y-takahashi\MyBrain\20_Projects\企業のオンライン保健室\クラウドパワー株式会社\02_報告会記録\クラウドパワー様_アンケート結果_2026年6月.pptx'
prs.save(out)
print(f'保存完了：{out}')
